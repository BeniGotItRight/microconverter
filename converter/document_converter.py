"""Document conversions: PDF, Word, Text, PPTX, HTML, Markdown, RTF, ODT."""

from pathlib import Path

from fpdf import FPDF


def _make_pdf() -> FPDF:
    """Create a FPDF instance with standard settings and Unicode font."""
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    # Use built-in Helvetica (latin-1 safe)
    pdf.set_font("Helvetica", size=11)
    return pdf


def _safe_latin1(text: str, max_len: int = 0) -> str:
    """Encode text to latin-1 safely for fpdf2."""
    if max_len:
        text = text[:max_len]
    return str(text).encode("latin-1", "replace").decode("latin-1")


def ocr_pdf(path: Path) -> str:
    """
    Extract text from scanned PDFs using OCR.
    Requires tesseract-ocr and poppler-utils system packages.
    """
    try:
        from pdf2image import convert_from_path
        import pytesseract
        
        # Check if tesseract is installed
        try:
            pytesseract.get_tesseract_version()
        except:
            raise RuntimeError("Tesseract OCR engine not found on system path.")
        
        images = convert_from_path(str(path), dpi=300)
        full_text = ""
        for i, img in enumerate(images):
            text = pytesseract.image_to_string(img, lang='eng')
            full_text += f"--- Page {i+1} ---\n{text}\n\n"
        return full_text
    except Exception as e:
        raise RuntimeError(
            "OCR failed. Please ensure 'tesseract-ocr' and 'poppler-utils' "
            f"are installed on your system. Error: {str(e)}"
        )


def pdf_to_word(input_path: Path, output_path: Path) -> None:
    """Convert PDF to Word (.docx) using pdf2docx."""
    from pdf2docx import Converter

    cv = Converter(str(input_path))
    try:
        cv.convert(str(output_path), start=0, end=None)
    finally:
        cv.close()


def pdf_to_text(input_path: Path, output_path: Path) -> None:
    """Extract all text from PDF and save as .txt. Falls back to OCR if needed."""
    import pdfplumber

    lines = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                lines.append(text)
    
    if not lines:
        # Try OCR as fallback
        try:
            text = ocr_pdf(input_path)
            lines = [text]
        except Exception as e:
            raise ValueError(f"No extractable text found and OCR failed: {str(e)}")
            
    output_path.write_text("\n\n".join(lines), encoding="utf-8")


def word_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert Word (.docx) to PDF using python-docx + fpdf2."""
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(str(input_path))
    pdf = _make_pdf()
    pdf.add_page()

    for block in doc.iter_inner_content():
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if text:
                pdf.multi_cell(0, 6, _safe_latin1(text))
        elif isinstance(block, Table):
            pdf.ln(4)
            for row in block.rows:
                cells = [str(cell.text or "").strip()[:80] for cell in row.cells]
                if any(cells):
                    line = " | ".join(cells)
                    pdf.multi_cell(0, 6, _safe_latin1(line))
            pdf.ln(4)
    pdf.output(str(output_path))


def word_to_text(input_path: Path, output_path: Path) -> None:
    """Extract text from Word (.docx) and save as .txt."""
    from docx import Document

    doc = Document(str(input_path))
    lines = []
    for para in doc.paragraphs:
        lines.append(para.text)
    output_path.write_text("\n".join(lines), encoding="utf-8")


def text_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert plain text (.txt) to PDF."""
    with open(input_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()

    pdf = _make_pdf()
    pdf.add_page()

    for line in content.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        pdf.multi_cell(0, 6, _safe_latin1(line))
    pdf.output(str(output_path))


def pptx_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert PowerPoint (.pptx) to PDF. Extracts text from all slides."""
    from pptx import Presentation

    prs = Presentation(str(input_path))
    pdf = _make_pdf()

    for i, slide in enumerate(prs.slides):
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 10, _safe_latin1(f"Slide {i + 1}"), new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", size=11)
        pdf.ln(4)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        pdf.multi_cell(0, 6, _safe_latin1(text))
                        pdf.ln(2)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip()[:60] for cell in row.cells]
                    line = " | ".join(cells)
                    pdf.multi_cell(0, 6, _safe_latin1(line))
                pdf.ln(4)

    if not prs.slides:
        raise ValueError("PowerPoint file has no slides.")
    pdf.output(str(output_path))


def html_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert HTML file to PDF. Strips tags, preserves text structure."""
    from bs4 import BeautifulSoup

    with open(input_path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    pdf = _make_pdf()
    pdf.add_page()

    # Extract title
    title = soup.find("title")
    if title and title.string:
        pdf.set_font("Helvetica", "B", 16)
        pdf.cell(0, 10, _safe_latin1(title.string.strip()), new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)
        pdf.set_font("Helvetica", size=11)

    # Process headings and paragraphs
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th", "pre", "blockquote"]):
        text = tag.get_text(strip=True)
        if not text:
            continue
        if tag.name in ("h1", "h2"):
            pdf.set_font("Helvetica", "B", 14)
            pdf.ln(4)
            pdf.multi_cell(0, 7, _safe_latin1(text))
            pdf.set_font("Helvetica", size=11)
        elif tag.name in ("h3", "h4", "h5", "h6"):
            pdf.set_font("Helvetica", "B", 12)
            pdf.ln(3)
            pdf.multi_cell(0, 6, _safe_latin1(text))
            pdf.set_font("Helvetica", size=11)
        elif tag.name == "li":
            pdf.multi_cell(0, 6, _safe_latin1(f"  • {text}"))
        else:
            pdf.multi_cell(0, 6, _safe_latin1(text))
        pdf.ln(1)

    pdf.output(str(output_path))


def markdown_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert Markdown to PDF via HTML intermediate."""
    import markdown as md
    import tempfile

    with open(input_path, "r", encoding="utf-8", errors="replace") as f:
        md_text = f.read()

    html_content = md.markdown(md_text, extensions=["tables", "fenced_code", "nl2br"])
    # Wrap in basic HTML
    full_html = f"<html><body>{html_content}</body></html>"

    # Write temp HTML, then convert to PDF
    tmp = Path(tempfile.mktemp(suffix=".html"))
    try:
        tmp.write_text(full_html, encoding="utf-8")
        html_to_pdf(tmp, output_path)
    finally:
        if tmp.exists():
            tmp.unlink()


def markdown_to_html(input_path: Path, output_path: Path) -> None:
    """Convert Markdown to styled HTML."""
    import markdown as md

    with open(input_path, "r", encoding="utf-8", errors="replace") as f:
        md_text = f.read()

    html_body = md.markdown(md_text, extensions=["tables", "fenced_code", "nl2br"])
    styled_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Converted Document</title>
<style>
body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
       max-width: 800px; margin: 2rem auto; padding: 0 1rem;
       line-height: 1.6; color: #1a1a1a; }}
h1, h2, h3 {{ color: #0f172a; }}
code {{ background: #f1f5f9; padding: 0.2em 0.4em; border-radius: 4px; font-size: 0.9em; }}
pre {{ background: #f1f5f9; padding: 1em; border-radius: 8px; overflow-x: auto; }}
table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
th, td {{ border: 1px solid #e2e8f0; padding: 0.5em 1em; text-align: left; }}
th {{ background: #f8fafc; font-weight: 600; }}
blockquote {{ border-left: 4px solid #10b981; margin: 1em 0; padding: 0.5em 1em; color: #475569; }}
</style>
</head>
<body>
{html_body}
</body>
</html>"""
    output_path.write_text(styled_html, encoding="utf-8")


def rtf_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert RTF to PDF. Strips RTF formatting, extracts text."""
    from striprtf.striprtf import rtf_to_text

    with open(input_path, "r", encoding="utf-8", errors="replace") as f:
        rtf_content = f.read()

    plain_text = rtf_to_text(rtf_content)
    if not plain_text.strip():
        raise ValueError("No text could be extracted from RTF file.")

    pdf = _make_pdf()
    pdf.add_page()
    for line in plain_text.split("\n"):
        pdf.multi_cell(0, 6, _safe_latin1(line))
    pdf.output(str(output_path))


def rtf_to_text(input_path: Path, output_path: Path) -> None:
    """Convert RTF to plain text."""
    from striprtf.striprtf import rtf_to_text as _rtf_to_text

    with open(input_path, "r", encoding="utf-8", errors="replace") as f:
        rtf_content = f.read()

    plain_text = _rtf_to_text(rtf_content)
    if not plain_text.strip():
        raise ValueError("No text could be extracted from RTF file.")
    output_path.write_text(plain_text, encoding="utf-8")


def odt_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert ODT (OpenDocument Text) to PDF."""
    from odf.opendocument import load
    from odf import text as odf_text
    from odf.text import P, H

    doc = load(str(input_path))
    pdf = _make_pdf()
    pdf.add_page()

    for element in doc.getElementsByType(P):
        txt = ""
        for node in element.childNodes:
            txt += str(node) if hasattr(node, '__str__') else ""
        txt = txt.strip()
        if txt:
            pdf.multi_cell(0, 6, _safe_latin1(txt))
            pdf.ln(1)

    for element in doc.getElementsByType(H):
        txt = ""
        for node in element.childNodes:
            txt += str(node) if hasattr(node, '__str__') else ""
        txt = txt.strip()
        if txt:
            pdf.set_font("Helvetica", "B", 13)
            pdf.multi_cell(0, 7, _safe_latin1(txt))
            pdf.set_font("Helvetica", size=11)
            pdf.ln(2)

    pdf.output(str(output_path))


def odt_to_text(input_path: Path, output_path: Path) -> None:
    """Convert ODT to plain text."""
    from odf.opendocument import load
    from odf.text import P, H

    doc = load(str(input_path))
    lines = []
    for element_type in (P, H):
        for element in doc.getElementsByType(element_type):
            txt = ""
            for node in element.childNodes:
                txt += str(node) if hasattr(node, '__str__') else ""
            txt = txt.strip()
            if txt:
                lines.append(txt)

    if not lines:
        raise ValueError("No text could be extracted from ODT file.")
    output_path.write_text("\n".join(lines), encoding="utf-8")
